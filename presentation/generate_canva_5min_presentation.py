from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = BASE_DIR / "presentation"
PROCESSED_DIR = BASE_DIR / "data" / "processed"
CHARTS_DIR = PROCESSED_DIR / "charts"
EXTRACT_DIR = BASE_DIR / "extract" / "data"

OUTPUT_PATH = PRESENTATION_DIR / "precious_metal_market_analysis_5min_canva.pptx"
NOTES_PATH = PRESENTATION_DIR / "precious_metal_market_analysis_5min_speaker_notes.md"

FLOWCHART_PATH = PRESENTATION_DIR / "gold_silver_analysis_flowchart.png"
EQUITY_CURVE_PATH = CHARTS_DIR / "equity_curve_comparison_all_strategies.png"

LAST_UPDATE_PATH = PROCESSED_DIR / "last_update.json"
RISK_PATH = PROCESSED_DIR / "recent_2year_risk_metrics.csv"
COMPARE_PATH = PROCESSED_DIR / "equity_curve_comparison_all_strategies.csv"
PREMIUM_PATH = PROCESSED_DIR / "premium_analysis.csv"

GOLD_FUTURES_PATH = EXTRACT_DIR / "gold_futures.csv"
SILVER_FUTURES_PATH = EXTRACT_DIR / "silver_futures.csv"
GOLD_SPOT_PATH = EXTRACT_DIR / "gold_spot.csv"
SILVER_SPOT_PATH = EXTRACT_DIR / "silver_spot.csv"

BG = RGBColor(247, 244, 238)
NAVY = RGBColor(23, 37, 84)
GOLD = RGBColor(184, 134, 11)
SLATE = RGBColor(71, 85, 105)
DARK = RGBColor(30, 41, 59)
WHITE = RGBColor(255, 255, 255)
SOFT = RGBColor(236, 240, 245)
GREEN = RGBColor(22, 101, 52)
RED = RGBColor(153, 27, 27)


def fmt_date(value: str) -> str:
    return pd.to_datetime(value).strftime("%Y-%m-%d")


def load_inputs() -> dict:
    risk_df = pd.read_csv(RISK_PATH)
    compare_df = pd.read_csv(COMPARE_PATH)
    premium_df = pd.read_csv(PREMIUM_PATH)
    last_update = json.load(open(LAST_UPDATE_PATH, "r", encoding="utf-8"))

    gold_futures = pd.read_csv(GOLD_FUTURES_PATH)
    silver_futures = pd.read_csv(SILVER_FUTURES_PATH)
    gold_spot = pd.read_csv(GOLD_SPOT_PATH)
    silver_spot = pd.read_csv(SILVER_SPOT_PATH)

    best = compare_df.sort_values("final_capital", ascending=False).iloc[0]
    gold_buy20 = risk_df[(risk_df["asset"] == "gold") & (risk_df["signal"] == "BUY") & (risk_df["holding_days"] == 20)].iloc[0]
    silver_buy20 = risk_df[(risk_df["asset"] == "silver") & (risk_df["signal"] == "BUY") & (risk_df["holding_days"] == 20)].iloc[0]
    latest_premium = premium_df.dropna(subset=["gold_zscore", "silver_zscore"]).iloc[-1]

    return {
        "risk_df": risk_df,
        "compare_df": compare_df,
        "premium_df": premium_df,
        "last_update": last_update,
        "best": best,
        "gold_buy20": gold_buy20,
        "silver_buy20": silver_buy20,
        "latest_premium": latest_premium,
        "gold_futures_range": (
            fmt_date(gold_futures["Date"].min()),
            fmt_date(gold_futures["Date"].max()),
            len(gold_futures),
        ),
        "silver_futures_range": (
            fmt_date(silver_futures["Date"].min()),
            fmt_date(silver_futures["Date"].max()),
            len(silver_futures),
        ),
        "gold_spot_range": (
            fmt_date(gold_spot["date"].min()),
            fmt_date(gold_spot["date"].max()),
            len(gold_spot),
        ),
        "silver_spot_range": (
            fmt_date(silver_spot["date"].min()),
            fmt_date(silver_spot["date"].max()),
            len(silver_spot),
        ),
        "premium_range": (
            fmt_date(premium_df["date"].min()),
            fmt_date(premium_df["date"].max()),
            len(premium_df),
        ),
    }


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = DARK,
    align=PP_ALIGN.LEFT,
    font_name: str = "Aptos",
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(
    slide,
    items: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: int = 18,
    color: RGBColor = DARK,
    bullet_color: RGBColor = GOLD,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.bullet = True
        if p.runs:
            p.runs[0].font.color.rgb = color
    return box


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, 0.65, 0.35, 8.8, 0.55, font_size=25, bold=True, color=NAVY)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(0.95), Inches(2.1), Inches(0.07))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.color.rgb = GOLD
    if subtitle:
        add_text(slide, subtitle, 0.65, 1.05, 11.6, 0.35, font_size=10, color=SLATE)


def add_panel(slide, left: float, top: float, width: float, height: float, *, fill_color: RGBColor = WHITE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = SOFT
    return shape


def add_metric_card(slide, title: str, value: str, note: str, left: float, top: float, width: float = 2.7) -> None:
    add_panel(slide, left, top, width, 1.45)
    add_text(slide, title, left + 0.18, top + 0.12, width - 0.36, 0.25, font_size=11, color=SLATE)
    add_text(slide, value, left + 0.18, top + 0.45, width - 0.36, 0.45, font_size=24, bold=True, color=NAVY)
    add_text(slide, note, left + 0.18, top + 0.94, width - 0.36, 0.25, font_size=10, color=SLATE)


def add_source_card(slide, title: str, items: list[str], left: float, top: float, width: float, accent: RGBColor) -> None:
    add_panel(slide, left, top, width, 2.15)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    add_text(slide, title, left + 0.18, top + 0.22, width - 0.36, 0.28, font_size=15, bold=True, color=NAVY)
    add_bullets(slide, items, left + 0.12, top + 0.55, width - 0.24, 1.4, font_size=12)


def add_formula_box(slide, title: str, formula: str, note: str, left: float, top: float, width: float, accent: RGBColor) -> None:
    add_panel(slide, left, top, width, 1.5, fill_color=WHITE)
    add_text(slide, title, left + 0.18, top + 0.12, width - 0.36, 0.24, font_size=12, bold=True, color=accent)
    add_text(slide, formula, left + 0.18, top + 0.45, width - 0.36, 0.35, font_size=18, bold=True, color=NAVY)
    add_text(slide, note, left + 0.18, top + 0.95, width - 0.36, 0.22, font_size=10, color=SLATE)


def add_picture(slide, path: Path, left: float, top: float, width: float, height: float | None = None) -> None:
    if not path.exists():
        return
    if height is None:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def build_presentation(data: dict) -> Presentation:
    prs = new_prs()

    last_update = data["last_update"]
    best = data["best"]
    gold_buy20 = data["gold_buy20"]
    silver_buy20 = data["silver_buy20"]
    latest_premium = data["latest_premium"]
    gold_f_start, gold_f_end, gold_f_rows = data["gold_futures_range"]
    silver_f_start, silver_f_end, silver_f_rows = data["silver_futures_range"]
    gold_s_start, gold_s_end, gold_s_rows = data["gold_spot_range"]
    silver_s_start, silver_s_end, silver_s_rows = data["silver_spot_range"]
    premium_start, premium_end, premium_rows = data["premium_range"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    top_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(5.95))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = NAVY
    top_band.line.color.rgb = NAVY
    bottom_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(5.95), Inches(13.333), Inches(1.55))
    bottom_band.fill.solid()
    bottom_band.fill.fore_color.rgb = GOLD
    bottom_band.line.color.rgb = GOLD
    add_text(slide, "貴金屬市場分析 ETL 專案", 0.82, 1.05, 9.2, 0.8, font_size=28, bold=True, color=WHITE)
    add_text(slide, "黃金與白銀期貨/現貨 Premium、Z-score 與交易訊號建構", 0.82, 1.95, 10.5, 0.45, font_size=17, color=WHITE)
    add_text(slide, "5 分鐘簡報版｜Canva 匯入友善格式", 0.82, 2.48, 6.8, 0.35, font_size=14, color=WHITE)
    add_text(slide, f"資料更新：{last_update['last_update_time']}", 0.82, 3.4, 4.3, 0.28, font_size=13, color=WHITE)
    add_text(slide, f"最佳策略：{best['asset'].title()} {best['signal']} {int(best['holding_days'])}D | {best['total_return_pct']:.2f}%", 0.82, 3.82, 6.8, 0.3, font_size=13, color=WHITE)
    add_text(slide, "Patric Xie", 0.82, 6.45, 2.4, 0.3, font_size=20, bold=True, color=NAVY)
    add_text(slide, "Data Engineering / Quant Project Showcase", 0.82, 6.8, 5.3, 0.25, font_size=11, color=NAVY)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "1. 專案背景與目標", "把市場資料分析題目，做成可重複執行的 ETL pipeline")
    add_panel(slide, 0.65, 1.55, 5.95, 4.95)
    add_panel(slide, 6.78, 1.55, 5.9, 4.95)
    add_text(slide, "專案背景", 0.92, 1.82, 2.2, 0.3, font_size=16, bold=True, color=NAVY)
    add_bullets(slide, [
        "黃金與白銀同時具備避險與商品屬性，價格偏離具研究價值。",
        "我希望不只做單次分析，而是完整實作 Extract / Transform / Load。",
        "主題聚焦在期貨相對現貨的 Premium，觀察是否出現可量化的偏離。",
    ], 0.86, 2.2, 5.25, 3.7, font_size=15)
    add_text(slide, "專案目標", 7.05, 1.82, 2.2, 0.3, font_size=16, bold=True, color=NAVY)
    add_bullets(slide, [
        "自動抓取黃金與白銀的期貨、現貨資料。",
        "計算 Premium 與 Z-score，產生 BUY / SELL / HOLD 訊號。",
        "把分析結果輸出成 CSV 並存進 SQLite，方便後續查詢與擴充。",
    ], 6.99, 2.2, 5.15, 3.7, font_size=15)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "2. 資料來源", "資料涵蓋期貨、現貨與對齊後的分析資料集")
    add_source_card(
        slide,
        "黃金期貨",
        [f"來源：Yahoo Finance（GC=F）", f"期間：{gold_f_start} ~ {gold_f_end}", f"筆數：{gold_f_rows} 筆，日頻資料"],
        0.68,
        1.55,
        3.0,
        GOLD,
    )
    add_source_card(
        slide,
        "白銀期貨",
        [f"來源：Yahoo Finance（SI=F）", f"期間：{silver_f_start} ~ {silver_f_end}", f"筆數：{silver_f_rows} 筆，日頻資料"],
        3.92,
        1.55,
        3.0,
        GOLD,
    )
    add_source_card(
        slide,
        "現貨資料",
        [f"來源：Alpha Vantage API", f"黃金：{gold_s_start} ~ {gold_s_end}", f"白銀：{silver_s_start} ~ {silver_s_end}"],
        7.16,
        1.55,
        2.95,
        NAVY,
    )
    add_source_card(
        slide,
        "對齊後分析集",
        [f"期間：{premium_start} ~ {premium_end}", f"共 {premium_rows} 筆對齊資料", "欄位：premium%、z-score、signal"],
        10.34,
        1.55,
        2.31,
        NAVY,
    )
    add_panel(slide, 0.68, 4.18, 11.97, 1.9, fill_color=WHITE)
    add_text(slide, "資料選擇理由", 0.95, 4.45, 2.2, 0.25, font_size=14, bold=True, color=NAVY)
    add_bullets(slide, [
        "期貨反映市場預期，現貨更接近即時基礎價格，兩者的價差可以轉成 Premium 指標。",
        "將不同來源資料做日期對齊後，就能進一步計算偏離程度並建立統計訊號。",
    ], 0.9, 4.8, 11.1, 1.0, font_size=14)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "3. ETL 架構", "流程拆成 extract、transform、load，再由主程式串接執行")
    add_picture(slide, FLOWCHART_PATH, 0.72, 1.55, 7.1)
    add_panel(slide, 8.2, 1.62, 4.45, 1.35)
    add_panel(slide, 8.2, 3.1, 4.45, 1.5)
    add_panel(slide, 8.2, 4.78, 4.45, 1.4)
    add_text(slide, "Extract", 8.45, 1.84, 1.3, 0.25, font_size=14, bold=True, color=NAVY)
    add_text(slide, "抓取期貨 CSV 與現貨 API 資料，整理成可合併格式。", 8.45, 2.15, 3.8, 0.45, font_size=12, color=DARK)
    add_text(slide, "Transform", 8.45, 3.34, 1.6, 0.25, font_size=14, bold=True, color=NAVY)
    add_text(slide, "標準化日期/欄位、合併資料源、計算 Premium、Z-score 與訊號。", 8.45, 3.65, 3.8, 0.55, font_size=12, color=DARK)
    add_text(slide, "Load", 8.45, 5.04, 1.1, 0.25, font_size=14, bold=True, color=NAVY)
    add_text(slide, "輸出 CSV、寫入 SQLite，讓分析與展示層可以直接使用。", 8.45, 5.35, 3.8, 0.42, font_size=12, color=DARK)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "4. 分析方法", "將期貨相對現貨的偏離轉成統計訊號")
    add_formula_box(
        slide,
        "Premium",
        "Premium (%) = (Futures - Spot) / Spot × 100",
        "觀察期貨相對現貨是溢價還是折價",
        0.72,
        1.6,
        5.9,
        GOLD,
    )
    add_formula_box(
        slide,
        "Rolling Z-score",
        "Z = (Premium - Mean20) / Std20",
        "用 20 日窗口判斷目前偏離是否落在極端區間",
        6.78,
        1.6,
        5.85,
        NAVY,
    )
    add_panel(slide, 0.72, 3.5, 5.9, 2.35)
    add_panel(slide, 6.78, 3.5, 5.85, 2.35)
    add_text(slide, "訊號規則", 0.95, 3.75, 1.5, 0.25, font_size=15, bold=True, color=NAVY)
    add_bullets(slide, [
        "Z-score < -2：判定偏低，給 BUY 訊號。",
        "Z-score > 2：判定偏高，給 SELL 訊號。",
        "其餘區間：維持 HOLD。",
    ], 0.9, 4.1, 5.0, 1.35, font_size=14)
    add_text(slide, "方法價值", 7.02, 3.75, 1.5, 0.25, font_size=15, bold=True, color=NAVY)
    add_bullets(slide, [
        "把主觀判斷轉成可解釋、可重複的量化規則。",
        "後續可直接接回測、報表、儀表板與排程自動化。",
        f"最新資料點：Gold Z {latest_premium['gold_zscore']:.2f} / Silver Z {latest_premium['silver_zscore']:.2f}",
    ], 6.96, 4.1, 5.1, 1.5, font_size=14)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "5. 遇到的問題與解法", "資料工程的關鍵不只在抓資料，更在資料品質與流程穩定性")
    add_panel(slide, 0.72, 1.6, 5.75, 4.95)
    add_panel(slide, 6.84, 1.6, 5.77, 4.95)
    add_text(slide, "主要問題", 0.98, 1.88, 1.7, 0.25, font_size=16, bold=True, color=NAVY)
    add_bullets(slide, [
        "期貨與現貨資料欄位名稱不一致。",
        "不同來源資料日期頻率不同，容易對不齊。",
        "API 抓取偶爾缺值或格式異常，需要額外清理。",
        "如果流程寫得太集中，後續擴充和除錯會變困難。",
    ], 0.92, 2.25, 4.95, 3.5, font_size=15)
    add_text(slide, "對應解法", 7.1, 1.88, 1.7, 0.25, font_size=16, bold=True, color=NAVY)
    add_bullets(slide, [
        "統一欄位命名與日期格式，先標準化再合併。",
        "採用 inner join 對齊日期，避免錯配造成假訊號。",
        "在 transform 階段先做缺值處理與 rolling 計算檢查。",
        "將流程拆成 extract / transform / load，提升維護性。",
    ], 7.04, 2.25, 4.95, 3.5, font_size=15)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "6. 成果展示", "專案已能把原始價格資料轉成可分析、可回測、可展示的結果")
    add_metric_card(slide, "最佳策略", "Silver BUY 20D", "近兩年回測最佳表現", 0.72, 1.55, 2.8)
    add_metric_card(slide, "總報酬", f"{silver_buy20['total_return_pct']:.2f}%", "白銀 BUY 20D", 3.72, 1.55, 2.35)
    add_metric_card(slide, "勝率", f"{silver_buy20['win_rate_pct']:.2f}%", "白銀 BUY 20D", 6.27, 1.55, 2.0)
    add_metric_card(slide, "黃金參考", f"{gold_buy20['total_return_pct']:.2f}%", "Gold BUY 20D", 8.47, 1.55, 2.1)
    add_metric_card(slide, "Max DD", f"{silver_buy20['max_drawdown_pct']:.2f}%", "白銀 BUY 20D", 10.77, 1.55, 1.84)
    add_picture(slide, EQUITY_CURVE_PATH, 0.78, 3.05, 7.25)
    add_panel(slide, 8.35, 3.15, 4.28, 2.9)
    add_text(slide, "可交付成果", 8.63, 3.42, 1.9, 0.25, font_size=15, bold=True, color=NAVY)
    add_bullets(slide, [
        "分析結果輸出成 CSV。",
        "整理後資料可寫入 SQLite。",
        "圖表、PDF 報告與簡報可同步更新。",
        "能作為作品集與面試展示素材。",
    ], 8.56, 3.78, 3.35, 1.8, font_size=14)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "7. 未來優化與結論", "從單次專題，進一步擴展成可持續維護的資料產品")
    add_panel(slide, 0.72, 1.62, 5.75, 4.85)
    add_panel(slide, 6.84, 1.62, 5.77, 4.85, fill_color=WHITE)
    add_text(slide, "下一步優化", 0.98, 1.92, 1.9, 0.25, font_size=16, bold=True, color=NAVY)
    add_bullets(slide, [
        "加入排程自動化，固定更新資料與報表。",
        "部署到雲端環境，讓 dashboard 可持續存取。",
        "增加策略回測、交易成本與更多因子比較。",
        "把成果做成更完整的視覺化儀表板。",
    ], 0.92, 2.28, 4.95, 3.3, font_size=15)
    add_text(slide, "結論", 7.1, 1.92, 1.0, 0.25, font_size=16, bold=True, color=NAVY)
    add_bullets(slide, [
        "這個專案讓我完整實作一次資料工程流程。",
        "我不只處理分析邏輯，也處理資料品質、流程拆分與落地儲存。",
        "它讓我更確認自己想往資料工程方向發展。",
    ], 7.04, 2.28, 4.95, 2.3, font_size=15)
    add_text(
        slide,
        "一句話總結：把原始市場資料，轉成可重複執行、可查詢、可展示的 ETL 成果。",
        7.08,
        5.18,
        4.65,
        0.75,
        font_size=14,
        bold=True,
        color=GREEN,
    )

    return prs


def build_notes(data: dict) -> str:
    best = data["best"]
    gold_buy20 = data["gold_buy20"]
    silver_buy20 = data["silver_buy20"]
    last_update = data["last_update"]
    premium_start, premium_end, premium_rows = data["premium_range"]
    latest_premium = data["latest_premium"]

    return f"""# 貴金屬市場分析 5 分鐘講稿

## Slide 1 封面
各位好，今天我要報告的專案是貴金屬市場分析 ETL 專案。這個專案主要分析黃金與白銀的期貨價格、現貨價格與兩者之間的 Premium，並利用 Z-score 建立 BUY、SELL、HOLD 訊號。這份簡報使用的資料更新到 {last_update['last_update_time']}。

## Slide 2 專案背景與目標
我會做這個題目，是因為我希望不只是做單純資料分析，而是實際建立一套完整的資料流程，也就是從資料擷取、資料清理、資料轉換，到資料儲存的 ETL pipeline。專案目標是把原始價格資料轉成可分析、可查詢的交易訊號資料集。

## Slide 3 資料來源
資料來源包含黃金期貨、白銀期貨，以及現貨 API 資料。完成日期對齊後，目前分析資料集共有 {premium_rows} 筆，期間是 {premium_start} 到 {premium_end}。這樣我就可以比較期貨相對現貨的偏離程度。

## Slide 4 ETL 架構
整個專案架構分成 extract、transform、load 三個部分。extract 負責抓資料，transform 進行欄位整理、日期對齊與指標計算，load 則把結果輸出成 CSV 並寫入 SQLite，方便後續查詢與延伸應用。

## Slide 5 分析方法
在 transform 階段，我先計算期貨相對現貨的 Premium，再用 20 日 rolling window 計算 Z-score，判斷目前偏離是否落在異常高或異常低的區間。當 Z-score 小於 -2 時給 BUY，當 Z-score 大於 2 時給 SELL，其餘則為 HOLD。最新一筆資料中，Gold Z-score 是 {latest_premium['gold_zscore']:.2f}，Silver Z-score 是 {latest_premium['silver_zscore']:.2f}。

## Slide 6 遇到的問題與解法
在實作過程中，我遇到的主要問題包括資料來源格式不一致、部分資料抓取錯誤，以及欄位命名與日期對齊問題。我的做法是先統一欄位與日期格式，再做 inner join 對齊資料，並把 extract、transform、load 模組拆開，讓流程更穩定也更容易維護。

## Slide 7 成果展示
這個專案的成果，是成功建立了一套可重複執行的貴金屬分析流程。除了 Premium 與 Z-score 訊號之外，我也把資料接到回測分析。近兩年最佳策略是 {best['asset'].title()} {best['signal']} {int(best['holding_days'])}D，總報酬 {best['total_return_pct']:.2f}%。其中 Silver BUY 20D 的勝率是 {silver_buy20['win_rate_pct']:.2f}%，Gold BUY 20D 的總報酬則是 {gold_buy20['total_return_pct']:.2f}%。

## Slide 8 未來優化與結論
未來我希望再加入排程自動化、雲端部署、視覺化儀表板與更完整的策略回測，讓專案更完整。總結來說，這個專案讓我從資料抓取、清理、轉換到儲存，完整實作了一次資料工程流程，也讓我更確認自己想往資料工程方向發展。謝謝大家。
"""


def main() -> None:
    data = load_inputs()
    prs = build_presentation(data)
    prs.save(OUTPUT_PATH)
    NOTES_PATH.write_text(build_notes(data), encoding="utf-8")
    print(f"Saved PPTX: {OUTPUT_PATH}")
    print(f"Saved notes: {NOTES_PATH}")


if __name__ == "__main__":
    main()
