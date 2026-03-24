from __future__ import annotations

from pathlib import Path
import json

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parents[1]
PROCESSED_DIR = BASE_DIR / "data" / "processed"
CHARTS_DIR = PROCESSED_DIR / "charts"
OUTPUT_DIR = BASE_DIR / "presentation"

PPTX_PATH = OUTPUT_DIR / "precious_metal_market_analysis_10min.pptx"
PPTX_SIMPLE_PATH = OUTPUT_DIR / "precious_metal_market_analysis_10min_simple.pptx"
SCRIPT_PATH = OUTPUT_DIR / "precious_metal_market_analysis_10min_speaker_notes.md"

SUMMARY_PATH = PROCESSED_DIR / "recent_2year_backtest_report_table.csv"
RISK_PATH = PROCESSED_DIR / "recent_2year_risk_metrics.csv"
COMPARE_PATH = PROCESSED_DIR / "equity_curve_comparison_all_strategies.csv"
LAST_UPDATE_PATH = PROCESSED_DIR / "last_update.json"


BG = RGBColor(248, 246, 242)
NAVY = RGBColor(23, 37, 84)
GOLD = RGBColor(184, 134, 11)
SLATE = RGBColor(71, 85, 105)
DARK = RGBColor(30, 41, 59)
WHITE = RGBColor(255, 255, 255)


def load_inputs() -> dict:
    summary_df = pd.read_csv(SUMMARY_PATH)
    risk_df = pd.read_csv(RISK_PATH)
    compare_df = pd.read_csv(COMPARE_PATH)
    last_update = json.load(open(LAST_UPDATE_PATH, "r", encoding="utf-8"))

    best = compare_df.sort_values("final_capital", ascending=False).iloc[0]
    best_risk = risk_df.sort_values("final_capital", ascending=False).iloc[0]
    gold_buy20 = risk_df[(risk_df["asset"] == "gold") & (risk_df["signal"] == "BUY") & (risk_df["holding_days"] == 20)].iloc[0]
    silver_buy20 = risk_df[(risk_df["asset"] == "silver") & (risk_df["signal"] == "BUY") & (risk_df["holding_days"] == 20)].iloc[0]
    sell_rows = risk_df[risk_df["signal"] == "SELL"].sort_values("total_return_pct")

    return {
        "summary_df": summary_df,
        "risk_df": risk_df,
        "compare_df": compare_df,
        "last_update": last_update,
        "best": best,
        "best_risk": best_risk,
        "gold_buy20": gold_buy20,
        "silver_buy20": silver_buy20,
        "worst_sell": sell_rows.iloc[0],
    }


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = NAVY

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.02), Inches(2.2), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.color.rgb = GOLD

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.12), Inches(11.8), Inches(0.4))
        p = sub_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(11)
        run.font.color.rgb = SLATE


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, font_size: int = 20) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)


def add_metric_card(slide, title: str, value: str, left: float, top: float, width: float = 2.1, height: float = 1.2) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(226, 232, 240)

    title_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.12), Inches(width - 0.3), Inches(0.3))
    p1 = title_box.text_frame.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(11)
    p1.font.color.rgb = SLATE

    value_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.45), Inches(width - 0.3), Inches(0.45))
    p2 = value_box.text_frame.paragraphs[0]
    p2.text = value
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = NAVY


def add_image(slide, path: Path, left: float, top: float, width: float, height: float | None = None) -> None:
    if height is None:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_table_text(slide, lines: list[str], left: float, top: float, width: float, height: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)


def build_presentation(data: dict) -> Presentation:
    prs = new_prs()

    best = data["best"]
    gold_buy20 = data["gold_buy20"]
    silver_buy20 = data["silver_buy20"]
    worst_sell = data["worst_sell"]
    last_update = data["last_update"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.color.rgb = NAVY
    gold_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(6.5), Inches(13.333), Inches(1.0))
    gold_band.fill.solid()
    gold_band.fill.fore_color.rgb = GOLD
    gold_band.line.color.rgb = GOLD

    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(10.8), Inches(1.5))
    p = title.text_frame.paragraphs[0]
    p.text = "貴金屬市場分析"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sub = title.text_frame.add_paragraph()
    sub.text = "黃金與白銀期貨的量化訊號、回測結果與自動化資料流程"
    sub.font.size = Pt(18)
    sub.font.color.rgb = WHITE

    info = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(8.0), Inches(1.2))
    p = info.text_frame.paragraphs[0]
    p.text = f"更新時間：{last_update['last_update_time']}"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p = info.text_frame.add_paragraph()
    p.text = f"最佳策略：{best['asset'].title()} {best['signal']} {int(best['holding_days'])}D | 總報酬 {best['total_return_pct']:.2f}%"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "1. 市場背景與研究動機", "為什麼選擇黃金與白銀，為什麼用量化方法做市場觀察")
    add_bullets(slide, [
        "黃金與白銀兼具避險、工業需求與美元利率敏感度，價格波動具有研究價值。",
        "貴金屬價格常受通膨預期、地緣政治、美元強弱與市場情緒共同影響。",
        "本專案的核心問題是：價格偏離均值後，是否存在可量化的均值回歸機會。",
        "研究目標不是預測所有行情，而是找出偏離過大時的高勝率交易區間。",
    ], 0.8, 1.6, 11.5, 4.8, 21)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "2. 資料來源與分析架構", "從資料擷取、轉換、載入，到分析與視覺化展示")
    add_bullets(slide, [
        "資料來源：Yahoo Finance，使用 gold futures（GC=F）與 silver futures（SI=F）。",
        "ETL 流程：extract -> transform -> load -> backtest -> chart/report -> Streamlit dashboard。",
        "技術堆疊：Python、Pandas、SQLite、Matplotlib、Streamlit、GitHub Actions。",
        "現在已完成每日自動更新，包含圖表與 PDF 報告，讓展示資料保持最新。",
    ], 0.8, 1.6, 6.0, 4.8, 18)
    add_metric_card(slide, "資料頻率", "Daily", 7.2, 1.8)
    add_metric_card(slide, "回測區間", "2 Years", 9.5, 1.8)
    add_metric_card(slide, "自動更新", "08:30", 7.2, 3.2)
    add_metric_card(slide, "前端展示", "Streamlit", 9.5, 3.2)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "3. 訊號邏輯與研究方法", "以 MA20 偏離程度與 rolling Z-score 建立 BUY / SELL / HOLD 訊號")
    add_bullets(slide, [
        "先計算 close 與 20 日均線的偏離：dev_ma20 = (close - ma_20) / ma_20。",
        "再用 60 日 rolling window 計算 dev_ma20 的 Z-score，衡量偏離程度是否異常。",
        "若 dev_ma20_z <= -1.5 且短期均線高於長期均線，視為順勢中的回檔，給 BUY。",
        "若 dev_ma20_z >= 1.5 且短期均線低於長期均線，視為弱勢中的反彈，給 SELL。",
        "最後分別測試 5、10、20 天持有期，觀察報酬、勝率、回撤與 Sharpe ratio。",
    ], 0.8, 1.55, 11.8, 5.2, 18)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "4. 核心回測結果", "結論很明確：BUY 策略整體優於 SELL，且白銀彈性最大")
    add_metric_card(slide, "最佳策略", "Silver BUY 20D", 0.9, 1.7, 2.5)
    add_metric_card(slide, "總報酬", f"{silver_buy20['total_return_pct']:.2f}%", 3.7, 1.7)
    add_metric_card(slide, "勝率", f"{silver_buy20['win_rate_pct']:.2f}%", 6.0, 1.7)
    add_metric_card(slide, "Max DD", f"{silver_buy20['max_drawdown_pct']:.2f}%", 8.3, 1.7)
    add_metric_card(slide, "Sharpe", f"{silver_buy20['sharpe_ratio']:.4f}", 10.6, 1.7)
    add_table_text(slide, [
        f"1. Silver BUY 20D: 資金 {silver_buy20['final_capital']:.2f}，總報酬 {silver_buy20['total_return_pct']:.2f}%。",
        f"2. Silver BUY 10D: 資金 {data['risk_df'].sort_values('final_capital', ascending=False).iloc[1]['final_capital']:.2f}，總報酬 111.29%。",
        f"3. Gold BUY 20D: 資金 {gold_buy20['final_capital']:.2f}，總報酬 {gold_buy20['total_return_pct']:.2f}%。",
        f"反向觀察：最差 SELL 策略為 {worst_sell['asset'].title()} SELL {int(worst_sell['holding_days'])}D，總報酬 {worst_sell['total_return_pct']:.2f}%。",
        "說明目前市場中，順勢回檔買進比逆勢放空更有效率。",
    ], 0.9, 3.25, 11.7, 2.7)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "5. 圖表解讀：策略績效比較", "資金曲線顯示白銀 BUY 20D 明顯領先，其次為白銀 BUY 10D 與黃金 BUY 20D")
    add_image(slide, CHARTS_DIR / "equity_curve_comparison_all_strategies.png", 0.8, 1.45, 7.5)
    add_bullets(slide, [
        "白銀的價格彈性較高，因此在均值回歸策略中，報酬擴張也最明顯。",
        "黃金表現相對穩健，總報酬不如白銀，但回撤更可控。",
        "資金曲線之間的差異，代表同樣的訊號框架在不同資產上有不同表現。",
    ], 8.55, 1.7, 4.0, 4.8, 17)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "6. 圖表解讀：訊號與風險", "Z-score 可以看出極端偏離點，風險指標則幫助判斷策略可執行性")
    add_image(slide, CHARTS_DIR / "gold_zscore_2year_signal.png", 0.75, 1.55, 5.9)
    add_image(slide, CHARTS_DIR / "silver_zscore_2year_signal.png", 6.75, 1.55, 5.9)
    note = slide.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(11.8), Inches(0.5))
    p = note.text_frame.paragraphs[0]
    p.text = f"風險比較：Gold BUY 20D Max Drawdown {gold_buy20['max_drawdown_pct']:.2f}% | Silver BUY 20D Max Drawdown {silver_buy20['max_drawdown_pct']:.2f}%"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "7. 系統價值與應用場景", "這不只是研究，而是一個可維護、可展示、可自動更新的量化分析系統")
    add_bullets(slide, [
        "資料工程面：建立 ETL、健康檢查、SQLite 落地、自動更新與版本管理。",
        "分析面：把原始價格轉成可操作的訊號、回測與風險指標。",
        "展示面：用 Streamlit 把結果做成 dashboard，適合面試、專題或作品集展示。",
        "商業應用：可延伸到投資研究、策略監控、商品市場觀察與自動化報表。",
    ], 0.8, 1.55, 11.7, 4.5, 20)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, "8. 結論、限制與下一步", "目前結果具參考價值，但仍需控制樣本偏誤與交易假設")
    add_bullets(slide, [
        "結論：近兩年最佳策略為 Silver BUY 20D，顯示白銀在極端偏離後的回歸效果最強。",
        "限制：目前未納入交易成本、滑價、槓桿限制與宏觀因子聯動。",
        "限制：回測樣本主要集中在近兩年，仍需擴展到更長時間區間驗證穩定性。",
        "下一步：加入 spot/futures premium、美元指數、利率、通膨或機器學習特徵做多因子比較。",
        "下一步：強化簡報與 dashboard，讓系統同時兼具研究與產品展示能力。",
    ], 0.8, 1.55, 11.7, 5.2, 18)

    return prs


def add_simple_bullets(slide, title: str, bullets: list[str]) -> None:
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for idx, bullet in enumerate(bullets):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)


def build_simple_presentation(data: dict) -> Presentation:
    prs = Presentation()
    best = data["best"]
    gold_buy20 = data["gold_buy20"]
    silver_buy20 = data["silver_buy20"]
    worst_sell = data["worst_sell"]
    last_update = data["last_update"]

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "貴金屬市場分析"
    slide.placeholders[1].text = (
        "黃金與白銀期貨的量化訊號、回測結果與自動化資料流程\n"
        f"更新時間：{last_update['last_update_time']}\n"
        f"最佳策略：{best['asset'].title()} {best['signal']} {int(best['holding_days'])}D | {best['total_return_pct']:.2f}%"
    )

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_simple_bullets(slide, "1. 市場背景與研究動機", [
        "黃金與白銀兼具避險與商品屬性，適合做策略比較。",
        "研究問題：價格偏離均值後，是否存在可量化的回歸機會。",
        "目標：找出高勝率且可重複驗證的交易規則。",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_simple_bullets(slide, "2. 資料來源與分析架構", [
        "資料來源：Yahoo Finance，使用 GC=F 與 SI=F。",
        "流程：ETL -> 特徵與訊號 -> SQLite -> 回測 -> Dashboard。",
        "技術：Python、Pandas、Matplotlib、Streamlit、GitHub Actions。",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_simple_bullets(slide, "3. 訊號邏輯與研究方法", [
        "計算 close 相對 MA20 的偏離程度。",
        "使用 60 日 rolling Z-score 衡量偏離是否極端。",
        "Z-score 低且趨勢向上時 BUY，高且趨勢向下時 SELL。",
        "比較 5、10、20 天持有期的績效與風險。",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_simple_bullets(slide, "4. 核心回測結果", [
        f"最佳策略：Silver BUY 20D，總報酬 {silver_buy20['total_return_pct']:.2f}%。",
        f"勝率 {silver_buy20['win_rate_pct']:.2f}%，最大回撤 {silver_buy20['max_drawdown_pct']:.2f}%。",
        f"Gold BUY 20D 也有 {gold_buy20['total_return_pct']:.2f}% 的總報酬。",
        f"SELL 策略相對較弱，最差為 {worst_sell['asset'].title()} SELL {int(worst_sell['holding_days'])}D。",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "5. 圖表解讀：策略績效比較"
    slide.shapes.add_picture(str(CHARTS_DIR / "equity_curve_comparison_all_strategies.png"), Inches(0.8), Inches(1.3), width=Inches(11.5))

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "6. 圖表解讀：訊號與風險"
    slide.shapes.add_picture(str(CHARTS_DIR / "gold_zscore_2year_signal.png"), Inches(0.3), Inches(1.3), width=Inches(6.1))
    slide.shapes.add_picture(str(CHARTS_DIR / "silver_zscore_2year_signal.png"), Inches(6.8), Inches(1.3), width=Inches(6.1))

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_simple_bullets(slide, "7. 系統價值與應用場景", [
        "資料工程：ETL、健康檢查、自動更新與版本管理。",
        "分析：從原始價格轉成可解釋的交易訊號與風險指標。",
        "展示：Streamlit dashboard、PDF 報告、簡報檔。",
    ])

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    add_simple_bullets(slide, "8. 結論、限制與下一步", [
        "結論：近兩年最佳策略為 Silver BUY 20D。",
        "限制：尚未納入交易成本、滑價與更長週期驗證。",
        "下一步：加入 premium、美元、利率與通膨等多因子。",
    ])

    return prs


def build_speaker_notes(data: dict) -> str:
    best = data["best"]
    gold_buy20 = data["gold_buy20"]
    silver_buy20 = data["silver_buy20"]
    worst_sell = data["worst_sell"]
    last_update = data["last_update"]
    return f"""# 貴金屬市場分析 10 分鐘講稿

## Slide 1 封面
大家好，今天我要分享的是我的貴金屬市場分析專案。這個專案聚焦在黃金與白銀期貨，目標是透過量化方法找出均值回歸的交易機會，並且把整個流程做成可自動更新的資料產品。這份簡報使用的資料最後更新時間是 {last_update['last_update_time']}，目前最佳策略是 {best['asset'].title()} {best['signal']} {int(best['holding_days'])} 天，總報酬是 {best['total_return_pct']:.2f}%。

## Slide 2 市場背景與研究動機
我選擇黃金與白銀，主要是因為這兩種資產同時具有避險屬性與商品屬性。黃金通常更受到美元、實質利率與避險需求影響，白銀除了金融屬性之外，也與工業需求有關，所以波動通常更大。對我來說，這代表它們很適合做策略比較。我想回答的問題是，當價格偏離均值太多時，是否能用系統化規則捕捉回歸機會，而不是只靠主觀判斷。

## Slide 3 資料來源與分析架構
在資料面，我用 Yahoo Finance 擷取黃金與白銀期貨資料。整個系統從 ETL 開始，先做 extract，再做 transform 產生特徵與訊號，接著 load 到 SQLite，最後再跑回測、風險指標、圖表與 PDF 報告。前端則是用 Streamlit 做成 dashboard。最近我也把 GitHub Actions 排程加進去，所以這套系統每天可以自動更新，對於展示專案完整度很有幫助。

## Slide 4 訊號邏輯與研究方法
訊號邏輯的核心，是觀察價格相對於 20 日均線的偏離程度。我先計算 close 與 MA20 的距離，再用 60 日 rolling window 計算這個偏離的 Z-score。如果 Z-score 很低，而且短期均線仍在長期均線之上，我把它視為順勢中的回檔，給 BUY 訊號。反過來，如果 Z-score 很高，且短期均線在長期均線之下，就視為弱勢反彈，給 SELL 訊號。最後我測試 5、10、20 天不同持有期，觀察哪一組在報酬與風險上最合理。

## Slide 5 核心回測結果
這頁是整份研究最重要的結果。最佳策略是 Silver BUY 20D，最終資金來到 {silver_buy20['final_capital']:.2f}，總報酬 {silver_buy20['total_return_pct']:.2f}%，勝率 {silver_buy20['win_rate_pct']:.2f}%，最大回撤 {silver_buy20['max_drawdown_pct']:.2f}%。第二名是 Silver BUY 10D，第三名是 Gold BUY 20D，總報酬 {gold_buy20['total_return_pct']:.2f}%。反過來看，SELL 策略整體偏弱，最差的是 {worst_sell['asset'].title()} SELL {int(worst_sell['holding_days'])} 天，這也說明在這段樣本期間，做順勢中的回檔買進，比逆勢放空更有效。

## Slide 6 圖表解讀：策略績效比較
這張圖是各策略的資金曲線比較。可以明顯看到白銀 BUY 20D 的資金成長最突出，代表白銀在極端偏離後的反彈幅度最大。不過黃金 BUY 20D 雖然報酬較低，卻呈現較穩健的曲線。這告訴我們，不同資產即使使用相同邏輯，策略表現與風險輪廓還是有差異，所以不能只看報酬，還要一起看可承受的波動。

## Slide 7 圖表解讀：訊號與風險
這一頁我把黃金與白銀的 Z-score 訊號圖放在一起。可以看到 BUY 與 SELL 訊號其實不是每天都出現，而是集中在明顯偏離的區間，這代表策略不是高頻交易，而是偏向事件型、機會型的進出場。從風險指標來看，Gold BUY 20D 的最大回撤是 {gold_buy20['max_drawdown_pct']:.2f}%，Silver BUY 20D 是 {silver_buy20['max_drawdown_pct']:.2f}%。所以白銀雖然報酬高，但承受的波動也更大。

## Slide 8 系統價值與應用場景
我認為這個專案的價值不只在研究結論，而在於它已經具備資料工程、分析和展示三個層次。第一是資料工程，因為我已經把 ETL、健康檢查、自動更新和版本管理串起來。第二是分析，因為它能從價格資料生成可解釋的策略訊號。第三是展示，因為我用 Streamlit 做成 dashboard，也能輸出 PDF 和簡報。這讓專案很適合作品集、面試展示，甚至未來擴充成研究工具。

## Slide 9 結論、限制與下一步
最後總結一下，目前這套方法在近兩年的結果中，最佳策略是 Silver BUY 20D，代表白銀在極端偏離後的均值回歸效果最強。不過這份研究仍有幾個限制，例如目前沒有納入交易成本、滑價和更長時間區間驗證，所以結果不能直接視為實盤保證。下一步，我會把 futures 與 spot premium、美元指數、利率與通膨等因子一起納入，做成多因子版本，讓策略更完整。以上是我的分享，謝謝大家。
"""


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    data = load_inputs()
    prs = build_presentation(data)
    prs.save(PPTX_PATH)
    simple_prs = build_simple_presentation(data)
    simple_prs.save(PPTX_SIMPLE_PATH)
    SCRIPT_PATH.write_text(build_speaker_notes(data), encoding="utf-8")
    print(f"Saved PPTX: {PPTX_PATH}")
    print(f"Saved simple PPTX: {PPTX_SIMPLE_PATH}")
    print(f"Saved notes: {SCRIPT_PATH}")


if __name__ == "__main__":
    main()
