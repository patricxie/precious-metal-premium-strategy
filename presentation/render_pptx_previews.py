from __future__ import annotations

from io import BytesIO
from pathlib import Path
import textwrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


BASE_DIR = Path(__file__).resolve().parent
PPTX_PATH = BASE_DIR / "precious_metal_market_analysis_10min_simple.pptx"
OUT_DIR = BASE_DIR / "previews_full"

WIDTH = 1600
HEIGHT = 900


def emu_to_px_x(value: int, slide_width: int) -> int:
    return round(value / slide_width * WIDTH)


def emu_to_px_y(value: int, slide_height: int) -> int:
    return round(value / slide_height * HEIGHT)


def get_font(size: int, bold: bool = False):
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Supplemental/Helvetica.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            try:
                return ImageFont.truetype(path, size=size)
            except Exception:
                continue
    return ImageFont.load_default()


FONT_TITLE = get_font(34, bold=True)
FONT_BODY = get_font(24)
FONT_SMALL = get_font(18)


def draw_textbox(draw: ImageDraw.ImageDraw, shape, slide_width: int, slide_height: int) -> None:
    if not shape.has_text_frame:
        return

    x = emu_to_px_x(shape.left, slide_width)
    y = emu_to_px_y(shape.top, slide_height)
    w = emu_to_px_x(shape.width, slide_width)
    h = emu_to_px_y(shape.height, slide_height)

    cursor_y = y
    for p in shape.text_frame.paragraphs:
        text = p.text.strip()
        if not text:
            continue

        font = FONT_TITLE if p.font.size and p.font.size.pt >= 20 else FONT_BODY
        if p.level > 0:
            font = FONT_SMALL

        approx_chars = max(10, int(w / max(10, font.size * 0.55)))
        lines = textwrap.wrap(text, width=approx_chars) or [text]
        for idx, line in enumerate(lines):
            prefix = "• " if p.level == 0 and cursor_y > y + 10 and idx == 0 and len(shape.text_frame.paragraphs) > 1 else ""
            draw.text((x, cursor_y), prefix + line, fill=(30, 41, 59), font=font)
            cursor_y += int(font.size * 1.45)


def paste_picture(canvas: Image.Image, shape, slide_width: int, slide_height: int) -> None:
    image = Image.open(BytesIO(shape.image.blob)).convert("RGB")
    x = emu_to_px_x(shape.left, slide_width)
    y = emu_to_px_y(shape.top, slide_height)
    w = emu_to_px_x(shape.width, slide_width)
    h = emu_to_px_y(shape.height, slide_height)
    image = image.resize((w, h))
    canvas.paste(image, (x, y))


def render() -> None:
    prs = Presentation(PPTX_PATH)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    for idx, slide in enumerate(prs.slides, start=1):
        canvas = Image.new("RGB", (WIDTH, HEIGHT), (255, 255, 255))
        draw = ImageDraw.Draw(canvas)

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                paste_picture(canvas, shape, slide_width, slide_height)

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                draw_textbox(draw, shape, slide_width, slide_height)

        out = OUT_DIR / f"slide_{idx:02d}.png"
        canvas.save(out)
        print(f"Saved {out}")


if __name__ == "__main__":
    render()
